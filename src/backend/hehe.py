from flask import Flask, render_template, request, jsonify
from flask_mysqldb import MySQL
from flask_cors import CORS
import json
from pptx import Presentation
from io import BytesIO
import os
import torch
import torchaudio


app = Flask(__name__)
cors = CORS(app)

# Configure MySQL
app.config['MYSQL_HOST'] = 'localhost'
app.config['MYSQL_USER'] = 'root'
app.config['MYSQL_PASSWORD'] = 'root'
app.config['MYSQL_DB'] = 'fyp1'

mysql = MySQL(app)

def execute_query(query, values=None):
    try:
        cur = mysql.connection.cursor()
        if values:
            cur.execute(query, values)
        else:
            cur.execute(query)

        result = cur.fetchall()
        mysql.connection.commit()
        cur.close()
        return "DB Action performed successfully"
    except Exception as e:  
        return f'Database connection error: {str(e)}'

# @app.route('/add_user')
# def add_user():
#     query1="select * from user"
#     query2="INSERT INTO user (username, name, password, email) VALUES ('afzoo ', 'shosho', '5677', 'khokho@gmail.com')"
#     # Try to connect to the database
#     execute_query(query2)
#     return json.loads("done")


@app.route('/voice_model')
def voice_model():
   
    import requests
    import os
    import numpy as np
    import time

    # Silero STT model
    device = torch.device('cpu')
    model, decoder, utils = torch.hub.load(repo_or_dir='snakers4/silero-models',
                                        model='silero_stt',
                                        language='en', # also available 'de', 'es'
                                        device=device)
    (read_batch, split_into_batches,
    read_audio, prepare_model_input) = utils

    # Pyaudio recording setup
    import pyaudio
    import wave

    CHUNK = 1024
    FORMAT = pyaudio.paInt16
    CHANNELS = 1
    RATE = 16000
    RECORD_SECONDS = 10
    WAVE_OUTPUT_FILENAME = "output.wav"

    p = pyaudio.PyAudio()
    stream = p.open(format=FORMAT,
                    channels=CHANNELS,
                    rate=RATE,
                    input=True,
                    frames_per_buffer=CHUNK)

    print("* recording")

    frames = []

    # Start recording
    for i in range(0, int(RATE / CHUNK * RECORD_SECONDS)):
        data = stream.read(CHUNK)
        frames.append(data)

    print("* done recording")

    # Stop recording and close the stream
    stream.stop_stream()
    stream.close()
    p.terminate()


    start_time = time.time()

    # Save the recording to a file
    wf = wave.open(WAVE_OUTPUT_FILENAME, 'wb')
    wf.setnchannels(CHANNELS)
    wf.setsampwidth(p.get_sample_size(FORMAT))
    wf.setframerate(RATE)
    wf.writeframes(b''.join(frames))
    wf.close()

    WAVE_OUTPUT_FILENAME = "C:/Users/ayar3/Downloads/FYP-last/updated.wav"


    # Load the recording and transcribe
    test_files = [WAVE_OUTPUT_FILENAME]
    batches = split_into_batches(test_files, batch_size=10)
    input = prepare_model_input(read_batch(batches[0]),
                                device=device)





    output = model(input)


    end_time = time.time()
    for example in output:
        transcription = decoder(example.cpu())
        print(transcription)

    print("text:", transcription)

    execution_time = end_time - start_time
    print(f"Execution time: {execution_time:.2f} seconds")
    
    query1 = "INSERT INTO temppresentationslide (transcription) VALUES ( %s);"
    value=[]
    value.append(transcription)
    status=execute_query(query1, value)
    string_data = f'{{"transcription": "{transcription}", "execution_time": {execution_time}, "status": "{status}"}}'
    print("Received JSON data:", string_data)


    return json.loads(string_data)



    

@app.route('/upload_pptx', methods=['POST'])
def upload_pptx():
    try:
        pptx_file = request.files['pptxFile']
        pptx_content = BytesIO(pptx_file.read())
        presentation = Presentation(pptx_content)

        # Start a database transaction
        # This assumes you have a function start_transaction() and commit_transaction()
        
        for slide_number, slide in enumerate(presentation.slides, start=1):
            print(f"Slide {slide_number}:")
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += shape.text + "\n"

            # Insert the slide text into the database
            query = "INSERT INTO tempslides (slideNumber, slideText) VALUES (%s, %s);"
            values = (slide_number, slide_text)

            try:
                execute_query(query, values)
            except Exception as e:
                # Log the error for debugging
                print(f"Error inserting slide {slide_number}: {str(e)}")

        # Commit the database transaction
     
        string_data = '{"transcription": "done"}'
        return json.loads(string_data)

    except Exception as e:
        # Rollback the database transaction in case of an error
     
        # Log the error for debugging
        print(f'An error occurred: {str(e)}')
        return {'error': f'An error occurred: {str(e)}'}
    


    



if __name__ == '__main__':
    app.run(debug=True)