from flask import Flask, render_template, request, jsonify
from flask_cors import CORS
import json
from pptx import Presentation
from io import BytesIO
import threading
import pyaudio
import torch
import time
import psycopg2
import wave
import numpy as np
from flask import Flask, jsonify
from flask_cors import CORS

from sentence_transformers import SentenceTransformer, util
import nltk
from nltk.tokenize import sent_tokenize


#for importing keyPhrases file
from nltk.tokenize import word_tokenize
from rake_nltk import Rake
from nltk.corpus import stopwords


nltk.download('punkt')
nltk.download('stopwords')

app = Flask(__name__)
cors = CORS(app, origins=["http://localhost:3000"])

# Configure POSTGRES
app.config['POSTGRES_HOST'] = 'ep-shy-wind-a1rxd3hf.ap-southeast-1.aws.neon.tech'
app.config['POSTGRES_USER'] = 'msheharyarqazi'
app.config['POSTGRES_PASSWORD'] = '4uv9AnrLePWG'
app.config['POSTGRES_DB'] = 'intellectslide'


recording_status = {'status': 'record'} #global variable to update recording status
frames=[]
model=0
decoder=0
utils=0
stream=0
start_time=0
previous_frames_count = 0
previous_sentence=""
global_transcription=""
slideCount=1
presentationId=20
slideContent=""
# Load the model and get utility functions
model, decoder, utils = torch.hub.load(repo_or_dir='snakers4/silero-models',
                                    model='silero_stt',
                                    language='en',  # also available 'de', 'es'
                                    device=torch.device('cpu'))

# Assign utility functions to global variables
read_batch, split_into_batches, read_audio, prepare_model_input = utils



connection = psycopg2.connect(
    host=app.config['POSTGRES_HOST'],
    user=app.config['POSTGRES_USER'],
    password=app.config['POSTGRES_PASSWORD'],
    dbname=app.config['POSTGRES_DB']
        )

def execute_query(query, values=None):
    try:
       
        
        cur = connection.cursor()
        if values:
            cur.execute(query, values)
        else:
            cur.execute(query)
       

        # result = cur.fetchall()
        connection.commit()
        cur.close()
        # connection.close()
        # print("connection closed")
        return "Database operation successfull"
    except Exception as e:  
        print(f'Database connection error: {str(e)}')
        return f'Database connection error: {str(e)}'
    


#----------------------------------------------------------------

def context_match(query, slideData):
    embedder = SentenceTransformer("all-MiniLM-L6-v2")

    #sentences = nltk.sent_tokenize(slideData)
    sentences = slideData.split("\n")
    #sentences = sent_tokenize(slideData)

    # print("*********This is slide data**************** \n\n", sentences)
    # print("Hello testing, thisis a test line.\n")

    # Adding sentences to pupulate the list
    corpus = []
    corpus.extend(sentences)
    corpus_embeddings = embedder.encode(corpus, convert_to_tensor=True)

    # Adding sentences to populate the list
    queries = query

    query_embedding = embedder.encode(queries, convert_to_tensor=True)

    # We use cosine-similarity and torch.topk 
    cos_scores = util.cos_sim(query_embedding, corpus_embeddings)[0]
    top_results = torch.topk(cos_scores, k=1)

    print("\n\n======================\n\n")
    print("Query:", queries)

    for score, idx in zip(top_results[0], top_results[1]):
        matchedSentence = corpus[idx]
        # print('The match sentence is : ', matchedSentence)

    return matchedSentence

#----------------------------------------------------------------




#----------------------------------------------------------------

def keyPhrase(trans):
        
        sentences =  [trans]

        filtered_corpus = []
        set(stopwords.words('english'))
        for sentence in sentences:
        
            words = word_tokenize(sentence)
            
            # Remove stop words
            filtered_words = [word for word in words if word.lower() not in stopwords.words('english')]
            filtered_sentence = ' '.join(filtered_words)
            filtered_corpus.append(filtered_sentence)


        text = ' '.join(filtered_corpus)
        r = Rake()
        r.extract_keywords_from_text(text)
        key_phrases = r.get_ranked_phrases()
        for i, phrase in enumerate(key_phrases[:15]):
            print(f"Key Phrase {i + 1}: {phrase}")
       
        return key_phrases

#----------------------------------------------------------------


    
def start_record():
    global model
    global decoder
    global utils
    global stream
    global frames
    stream=0

    # Pyaudio recording setup
    CHUNK = 1024
    FORMAT = pyaudio.paInt16
    CHANNELS = 1
    RATE = 16000
    RECORD_SECONDS = 10
    WAVE_OUTPUT_FILENAME = "output.wav"

    p = pyaudio.PyAudio()
    stream = p.open(format=FORMAT,
                    channels=CHANNELS,
                    rate=RATE,
                    input=True,
                    frames_per_buffer=CHUNK)

    print("* recording")
    
    global start_time
    start_time = time.time()

    frames.clear()
    # Start recording
    # for i in range(0, int(RATE / CHUNK * RECORD_SECONDS)):
    global recording_status
    while recording_status['status'] == "record":
        data = stream.read(CHUNK)
        frames.append(data)
        
    


def stop_recording():

    global model
    global decoder
    global utils
    global stream
    global frames
    
    print("* done recording")

    # Stop recording and close the stream
    p = pyaudio.PyAudio()

    stream.stop_stream()
    stream.close()
    p.terminate()

    import wave
    CHUNK = 1024
    FORMAT = pyaudio.paInt16
    CHANNELS = 1
    RATE = 16000
    WAVE_OUTPUT_FILENAME = "output.wav"
    # Save the recording to a file
    wf = wave.open(WAVE_OUTPUT_FILENAME, 'wb')
    wf.setnchannels(CHANNELS)
    wf.setsampwidth(p.get_sample_size(FORMAT))
    wf.setframerate(RATE)
    wf.writeframes(b''.join(frames))
    wf.close()

    # Load the recording and transcribe
    test_files = [WAVE_OUTPUT_FILENAME]
    batches = split_into_batches(test_files, batch_size=10)
    input = prepare_model_input(read_batch(batches[0]),
                                device=torch.device('cpu'))

    WAVE_OUTPUT_FILENAME = "D:/University/FYP/FYP Repo/updated.wav"
    output = model(input)

    end_time = time.time()
    transcription=""
    for example in output:
        transcription = decoder(example.cpu())
        print(transcription)

    global start_time
    execution_time = end_time - start_time
    print(f"Execution time: {execution_time:.2f} seconds")
    
    #-----------------------------------------------------------
    #key Phrases

    
    key_phrases = keyPhrase(transcription)
    
    # query1 = "INSERT INTO presentationslide (transcription) VALUES ( %s);"
    query1 = "INSERT INTO presentationslide (transcription, keypoints) VALUES (%s, %s);"
    # value=[]
    values = [transcription, key_phrases]
    # value.append(transcription)
    # status=execute_query(query1, value)
    status = execute_query(query1, values)
    #-----------------------------------------------------------
    
    string_data = f'{{"transcription": "{transcription}", "execution_time": {execution_time:.2f}, "status": "{status}"}}'
    print("Received JSON data:", string_data)
    

    return json.loads(string_data)


# ----------------
def getSlideText():
    global slideCount
    global presentationId
    global slideContent
    print("Presentation Id: ", presentationId)
    print("Slide no: ", slideCount)
    query_getSlideContent = f"SELECT \"textContent\" FROM \"slide\" WHERE \"slideNo\" = {slideCount} and \"presentationId\" = {presentationId}"


    cur = connection.cursor()
    cur.execute(query_getSlideContent)
    slideContent = cur.fetchone()[0]
    connection.commit()
    cur.close()

    print("\nThis slide data\n")
    print(slideContent)

    return slideContent

#--------------------









def transcribe_data():
    global model
    global decoder
    global utils
    global stream
    global frames
    global previous_frames_count
    import copy
    # tempFrames = copy.copy(frames)
    tempFrames = copy.copy(frames[previous_frames_count:])

    # Update the previous_frames_count to the current count
    previous_frames_count = len(frames)

    
 
    p = pyaudio.PyAudio()

  

    import wave
    CHUNK = 1024
    FORMAT = pyaudio.paInt16
    CHANNELS = 1
    RATE = 16000
    WAVE_OUTPUT_FILENAME = "temp.wav"
    # Save the recording to a file
    wf = wave.open(WAVE_OUTPUT_FILENAME, 'wb')
    wf.setnchannels(CHANNELS)
    wf.setsampwidth(p.get_sample_size(FORMAT))
    wf.setframerate(RATE)
    wf.writeframes(b''.join(tempFrames))
    wf.close()

    # Load the recording and transcribe
    test_files = [WAVE_OUTPUT_FILENAME]
    batches = split_into_batches(test_files, batch_size=10)
    input = prepare_model_input(read_batch(batches[0]),
                                device=torch.device('cpu'))

    output = model(input)

    
    transcription=""
    for example in output:
        transcription = decoder(example.cpu())
        # print(transcription)

    return transcription



@app.route('/match_context', methods=['GET'])
def initialize_match_context():

    ST = time.time()
    global slideContent
    global global_transcription
    global previous_sentence
    previous_sentence=global_transcription
    global_transcription=[]
    global_transcription=transcribe_data()
    print(global_transcription)
    if(global_transcription is None):
        ET= time.time()
        ExecT= ET - ST
        return f'{{"matched sentence": "","execution_time": {ExecT:.2f}}}'


    length = len(previous_sentence)
    later_half = previous_sentence[length // 2:]

    # Concatenate the later half with the second string
    global_transcription = later_half + " " + global_transcription
    print(global_transcription)
    
    
    # ---- context Match func call -----
    matchedSentence = context_match(global_transcription, slideContent)


    print("\nThe matched sentence:\n")
    print(matchedSentence)


    ET= time.time()
    ExecT= ET - ST

    string_data = f'{{"matched sentence": "{matchedSentence}","execution_time": {ExecT:.2f}}}'
    return json.loads(string_data)


@app.route('/update_slide_count', methods=['POST'])
def update_slide_count():
    
    global slideCount
    global previous_frames_count
    global previous_sentence
    previous_sentence=""
    previous_frames_count = len(frames)
    operation = request.form.get('operation', 'next')
    if(operation == "next"):
        slideCount=slideCount + 1
    elif(operation == "previous"):
        slideCount=slideCount - 1
    getSlideText()
    result = f'{{"updated slide number": "{slideCount}"}}'

    return json.loads(result)



@app.route('/add_user', methods=['GET'])
def add_user():
    query2 = "INSERT INTO \"user\"  VALUES ('afzoo', 'shosho', '5677', 'khokho@gmail.com')"

    # Try to connect to the database
    execute_query(query2)
    

    string_data = f'{{"status": "user added"}}'

    return json.loads(string_data)


@app.route('/start_recording', methods=['GET'])
def start_model():
    global recording_status
    recording_status['status'] = "record"
    #The function below will recprd the whole session on a new thread
    threading.Thread(target=start_record).start()
    #The function below will match context periodically
    threading.Thread(target=call_function_periodically).start()
    string_data = f'{{"status": "recording started"}}'
    return json.loads(string_data)

# -----------------
def call_function_periodically():
    global previous_frames_count
    global frames
    global previous_sentence
    while recording_status['status'] == "record":
        diff=len(frames) - previous_frames_count
        if (diff > 100):
            print("----------------------Context Automatically Matched-------------------------------")
            initialize_match_context()
            
        time.sleep(5)  # Sleep for 5 seconds before calling the function again

# -------------------



@app.route('/stop_recording', methods=['GET'])
def stop_model():
    global recording_status
    recording_status['status'] = "stop_recording"
    return stop_recording()

from datetime import datetime


@app.route('/upload_pptx', methods=['POST'])
def upload_pptx():
    global presentationId

    try:
        pptx_file = request.files['pptxFile']   

        # Read the content of the file
        pptx_content = BytesIO(pptx_file.read())
        pptx_data = pptx_file.read()

        # Extracting title and creation date
        presentation = Presentation(pptx_content)
        title = presentation.core_properties.title

         # Truncate title if it exceeds 45 characters
        truncated_title = title[:45]

        # Generating the upload date
        upload_date = datetime.now().strftime('%Y-%m-%d %H:%M:%S')

        # Define the SQL query to insert the PPTX file data along with title and creation date
        query = "INSERT INTO \"presentation\" (\"pptxFile\", \"title\", \"uploadDate\") VALUES (%s, %s, %s);"

        # Execute the SQL query to insert the PPTX file data
        execute_query(query, (psycopg2.Binary(pptx_data), truncated_title, upload_date))

        # Retrieve the presentation ID after insertion
        query = "SELECT MAX(\"presentationId\") FROM \"presentation\";"
        cur = connection.cursor()
        cur.execute(query)
        presentationId = cur.fetchone()[0]
        connection.commit()
        cur.close()

        # Iterate over slides to extract text content
        for slide_number, slide in enumerate(presentation.slides, start=1):
            query1 = "INSERT INTO slide (\"textContent\", \"slideNo\", \"presentationId\")  VALUES (%s, %s, %s);"
            value = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    value = value + shape.text
            execute_query(query1, (value, slide_number, presentationId))

        string_data = json.dumps({"transcription": "done"})
        return string_data
    except Exception as e:
        return {'error': f'An error occurred: {str(e)}'}    

CORS(app)

@app.route('/presentation_data', methods=['GET'])
def get_presentation_data():
    try:
        # Define the SQL query to select presentation data
        query = "SELECT title, \"uploadDate\" FROM presentation WHERE title IS NOT NULL AND title <> '' AND \"uploadDate\"  IS NOT NULL;"



        cur = connection.cursor()
        cur.execute(query)
        presentation_data = cur.fetchall()
        cur.close()

        # Prepare the response
        response_data = []
        for row in presentation_data:
            title= row[0]
            dateUpload = row[1].strftime("%d %B %Y")  # Format date as desired

            # Convert pptx_file to base64 or any other appropriate format if needed
            response_data.append({
                # 'presentationId': presentation_id,
                'title': title,
                'uploadDate': dateUpload
                # 'uploadDate': upload_date
                # Add other fields as needed
            })

        # Print the response data
        print("Presentation data:", response_data)

        # Return the response
        return jsonify(response_data)

    except Exception as e:
        # Print the error message
        print(f'An error occurred: {str(e)}')

        # Return an empty list if no data is fetched
        return jsonify([])
    
    finally:
        # Close the cursor and connection
        if 'cur' in locals():
            cur.close()
        if 'connection' in locals():
            connection.close()









@app.route('/recent_slide_data', methods=['GET'])
def get_recent_slide_data():
    try:
        # Define the SQL query to select slide data for the specified session number
        query = "SELECT keypoints, transcription FROM presentationslide ORDER BY new_sessionid DESC LIMIT 1;"
        
        cur = connection.cursor()
        cur.execute(query)
        slide_data = cur.fetchall()  # Fetch all slide data
        cur.close()

        # List to store formatted slide data
        formatted_slide_data = []

        for slide in slide_data:
            # Extract keypoints and transcription from the slide
            keypoints_str = slide[0]
            transcription = slide[1]

            # Preprocess the keypoints string
            keypoints_list = [keypoints_str.strip('"{}')]

            # Format the slide data as a dictionary
            slide_dict = {
                'keyPoints': keypoints_list,
                'transcription': transcription
            }

            # Append the formatted slide data to the list
            formatted_slide_data.append(slide_dict)

        # Print the slide data on the console
        print("Slide data:", formatted_slide_data)

        # Return the formatted slide data as JSON
        return jsonify(formatted_slide_data)

    except Exception as e:
        # Print the error message
        print(f'An error occurred: {str(e)}')

        # Return an error response
        return jsonify({'error': 'Failed to retrieve slide data'})

    finally:
        # Close the connection
        if 'connection' in locals():
            connection.close()












if __name__ == '__main__':
   app.run(debug=True)